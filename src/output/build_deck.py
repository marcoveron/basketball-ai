"""Generate the project presentation (.pptx) with speaker notes.

A self-contained, professional 16:9 deck in the dashboard's dark theme. It reuses the
real artefacts the pipeline produced (game-flow chart, shot-distribution map, shot-mix
chart, player thumbnails) so the slides show actual results, not mock-ups. Every slide
carries presenter notes.

Usage:  python -m src.output.build_deck --runs runs720 --out Basketball_AI_Presentation.pptx
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# palette (matches the dashboard)
BG   = RGBColor(0x0B, 0x0D, 0x12)
PANEL= RGBColor(0x16, 0x19, 0x22)
INK  = RGBColor(0xEE, 0xF0, 0xF4)
MUT  = RGBColor(0x96, 0x9C, 0xAB)
ACC  = RGBColor(0x22, 0xC5, 0x53)
HOME = RGBColor(0xFF, 0x5D, 0xB1)
AWAY = RGBColor(0x4D, 0xA3, 0xFF)
GOLD = RGBColor(0xE0, 0xA8, 0x00)
LINE = RGBColor(0x2A, 0x2F, 0x3A)
FONT = "Arial"

EMU_IN = 914400


def _solid(shape, rgb):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def txt(s, text, l, t, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
        italic=False, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, kicker, title, color=ACC):
    txt(s, kicker.upper(), 0.62, 0.42, 11.5, 0.4, 13, color, bold=True)
    txt(s, title, 0.6, 0.74, 12.1, 1.1, 30, INK, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.6), Inches(0.62), Inches(0.06))
    _solid(bar, color); bar.line.fill.background(); bar.shadow.inherit = False


def panel(s, l, t, w, h, fill=PANEL, line=LINE, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    _solid(shp, fill)
    shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except (IndexError, KeyError):
        pass
    return shp


def panel_text(shp, lines, pad=0.16):
    """lines: list of (text, size, color, bold[, italic])."""
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(pad); tf.margin_right = Inches(pad)
    tf.margin_top = Inches(pad); tf.margin_bottom = Inches(pad)
    for i, ln in enumerate(lines):
        text, size, color, bold = ln[0], ln[1], ln[2], ln[3]
        italic = ln[4] if len(ln) > 4 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT          # autoshape's 1st paragraph defaults to centre
        p.space_after = Pt(4); p.line_spacing = 1.05
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT


def bullets(s, items, l, t, w, h, size=16, gap=8, color=INK, dot="—  "):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        text, c, bold = (it if isinstance(it, tuple) else (it, color, False))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        r = p.add_run(); r.text = dot + text
        r.font.size = Pt(size); r.font.color.rgb = c; r.font.bold = bold; r.font.name = FONT
    return tb


def pic(s, path, l, t, w=None, h=None):
    path = Path(path)
    if not path.exists():
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    return s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h)), w, h


def pic_fit(s, path, l, t, box_w, box_h, frame=True):
    """Place an image fitted (contain) inside a box, centred, optional panel frame."""
    path = Path(path)
    if not path.exists():
        return
    if frame:
        panel(s, l, t, box_w, box_h, fill=RGBColor(0x0C, 0x0E, 0x13))
    iw, ih = Image.open(path).size
    ar = iw / ih
    pad = 0.12
    bw, bh = box_w - 2 * pad, box_h - 2 * pad
    if bw / bh > ar:
        h = bh; w = h * ar
    else:
        w = bw; h = w / ar
    cl = l + (box_w - w) / 2; ct = t + (box_h - h) / 2
    s.shapes.add_picture(str(path), Inches(cl), Inches(ct), Inches(w), Inches(h))


def footer(s, idx, total):
    txt(s, "Basketball AI", 0.6, 7.04, 4, 0.3, 9, MUT)
    txt(s, f"{idx} / {total}", 12.1, 7.04, 0.7, 0.3, 9, MUT, align=PP_ALIGN.RIGHT)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def chip(s, label, l, t, color, w=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.34))
    _solid(shp, color); shp.line.fill.background(); shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RGBColor(0x06, 0x21, 0x0F)
    r.font.name = FONT


# --------------------------------------------------------------------------- slides
def build(R: Path, out: Path, repo_url: str = "github.com/marcoveron/basketball-ai") -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    TOTAL = 17
    n = 0

    # 1 — TITLE -------------------------------------------------------------
    s = slide(prs); n += 1
    # accent band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    _solid(band, ACC); band.line.fill.background(); band.shadow.inherit = False
    txt(s, "🏀  BASKETBALL AI", 0.9, 2.0, 11, 0.6, 18, ACC, bold=True)
    txt(s, "From raw game video to player analytics", 0.86, 2.55, 11.6, 1.3, 44, INK, bold=True)
    txt(s, "Amateur Basketball Player Personal-Data Intelligent Analysis System",
        0.9, 3.95, 11.3, 0.6, 18, MUT)
    txt(s, "Case study: T3BA Taiwan 3×3 championship final · 720p broadcast footage",
        0.9, 4.45, 11.3, 0.5, 14, MUT, italic=True)
    chip(s, "PROOF OF CONCEPT", 0.9, 5.25, ACC, w=2.1)
    notes(s, """
Welcome. This project is an AI system that turns ordinary amateur basketball footage into
the kind of structured analytics that, until now, only professional teams could afford.
Our test case is the final of a Taiwan 3x3 tournament — a single 720p broadcast clip, no
sensors, no manual tagging. Everything you'll see was produced automatically from that one
video file. I'll walk through the problem, how the system works, the results, and — just as
importantly — where we were honest about what this footage can and can't support.
""")
    footer(s, n, TOTAL)

    # 2 — PROBLEM -----------------------------------------------------------
    s = slide(prs); n += 1
    header(s, "The problem", "Amateur players fly blind")
    bullets(s, [
        ("Pro teams have player tracking, shot charts and efficiency metrics. Amateur games have a final score — and nothing else.", INK, False),
        ("The data exists in the video: every shot, every position, every player. It's just locked inside raw pixels.", INK, False),
        ("Manual tagging is slow, expensive and doesn't scale to weekend games.", INK, False),
    ], 0.62, 2.05, 7.0, 3.5, size=18, gap=14)
    panel(s, 8.05, 2.0, 4.6, 4.1, fill=RGBColor(0x12, 0x17, 0x1f))
    pt = s.shapes[-1]
    panel_text(pt, [
        ("THE GOAL", 13, ACC, True),
        ("", 6, INK, False),
        ("Input", 13, MUT, True),
        ("One raw game video", 17, INK, True),
        ("", 6, INK, False),
        ("Output", 13, MUT, True),
        ("Per-player stats · shot maps · auto-highlights · a shareable dashboard", 17, INK, True),
    ])
    notes(s, """
Here's the gap. A professional team walks off the court with shot charts, efficiency
numbers and tracking data. An amateur team walks off with the score on a phone photo.
But all of that information is already present in the footage — every shot and every
movement is there in the pixels. The only thing missing is something to extract it
automatically. Manual video tagging exists, but it's slow and expensive and nobody does it
for a Tuesday-night pickup game. Our goal, on the right: take one raw video in, and give
back professional-grade outputs.
""")
    footer(s, n, TOTAL)

    # 3 — WHAT WE BUILT -----------------------------------------------------
    s = slide(prs); n += 1
    header(s, "The solution", "One automated pipeline")
    txt(s, "Raw .mp4  →  detect  →  identify  →  read the game  →  dashboard, cards & highlights",
        0.62, 1.95, 12, 0.5, 16, MUT, italic=True)
    cards = [
        ("DETECTION", "Find the ball, basket & every\nplayer in each frame", HOME),
        ("IDENTITY", "Read jersey numbers (OCR)\n& track players across frames", AWAY),
        ("THE GAME", "Read the scoreboard to recover\nevery made shot, value & time", GOLD),
        ("OUTPUT", "Dashboard, player cards,\nshot map & auto-cut highlights", ACC),
    ]
    x = 0.62; w = 2.95; gap = 0.18
    for i, (k, body, c) in enumerate(cards):
        p = panel(s, x + i * (w + gap), 2.5, w, 3.2)
        panel_text(p, [(k, 14, c, True), ("", 8, INK, False), (body, 15, INK, False)])
    notes(s, """
Our answer is a single pipeline. You feed it a raw mp4 and it runs four stages end to end.
First, detection: a trained model finds the ball, the basket and every player in every
frame. Second, identity: we read jersey numbers with OCR and track each player across the
game. Third — and this is the clever part I'll come back to — we read the scoreboard itself
to recover every made shot. And fourth, output: it assembles a dashboard, player cards, a
shot-distribution map and automatically cut highlight clips. No human in the loop after you
press go.
""")
    footer(s, n, TOTAL)

    # 4 — ARCHITECTURE (4 layers) ------------------------------------------
    s = slide(prs); n += 1
    header(s, "Architecture", "Four functional layers")
    layers = [
        ("1 · COURT MODELING & DETECTION", "Detect hoop, backboard and players; establish the spatial reference for everything downstream.", HOME),
        ("2 · PLAYER TRACKING & IDENTITY", "Jersey-number OCR as the primary ID, tracking to stitch fragments; facial recognition as fallback.", AWAY),
        ("3 · BEHAVIOR RECOGNITION", "Detect shooting events (value + timing) and rebounding; classify inside vs outside the arc.", GOLD),
        ("4 · OUTPUT GENERATION", "Shot heat maps, growth curves, auto-edited highlight clips and the analytics dashboard.", ACC),
    ]
    y = 2.0
    for k, body, c in layers:
        p = panel(s, 0.62, y, 12.1, 1.16)
        barl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(y), Inches(0.1), Inches(1.16))
        _solid(barl, c); barl.line.fill.background(); barl.shadow.inherit = False
        panel_text(p, [(k, 14, c, True), (body, 15, INK, False)], pad=0.18)
        y += 1.28
    notes(s, """
Architecturally the system is four layers, each building on the one below. Layer one is
court modeling and detection — finding the hoop and the players and setting up a spatial
frame of reference. Layer two is identity — who is who — driven mainly by reading jersey
numbers, with face recognition as a backup. Layer three is behavior recognition — turning
raw positions into basketball events: shots and rebounds. And layer four is output — the
heat maps, growth curves, highlight clips and dashboard that a player actually sees. This
mirrors the original product spec one-to-one.
""")
    footer(s, n, TOTAL)

    # 5 — TECH / PIPELINE ---------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Under the hood", "The processing pipeline")
    steps = ["detect", "scene-filter", "jersey OCR", "scoreboard OCR", "shot events", "player cards", "highlights", "dashboard"]
    x = 0.62
    for i, st in enumerate(steps):
        w = 1.42
        p = panel(s, x, 2.1, w, 0.7, fill=RGBColor(0x12, 0x17, 0x1f))
        tf = p.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = st
        rr.font.size = Pt(11.5); rr.font.bold = True; rr.font.color.rgb = INK; rr.font.name = FONT
        x += w + 0.085
    bullets(s, [
        ("Detection — YOLO model (Roboflow-trained, classes: ball / basket / person) running on the full 720p frames.", INK, False),
        ("Scene filter — drops replays & close-ups, keeping the ~81% of frames shot from the main game camera.", INK, False),
        ("Jersey OCR & scoreboard OCR — PaddleOCR on GPU, with per-region crops tuned for tiny on-court numbers.", INK, False),
        ("Highlights — ffmpeg cuts & captions one clip per made shot, then stitches a full reel.", INK, False),
    ], 0.62, 3.2, 12.1, 3.2, size=16, gap=12)
    notes(s, """
A quick look under the hood for the technical audience. Detection uses a YOLO model we
trained on labelled basketball data, with three classes — ball, basket and person — run on
the full-resolution frames. A scene filter throws away replays and close-ups and keeps the
roughly eighty percent of frames coming from the main game camera. Both the jersey numbers
and the scoreboard are read with PaddleOCR on the GPU, using tight per-region crops because
the numbers are tiny. And the highlights are cut and captioned with ffmpeg. Off-the-shelf
components, but the way they're wired together is what makes it work on messy amateur
footage.
""")
    footer(s, n, TOTAL)

    # 6 — ML: DETECTION MODEL / FINE-TUNING --------------------------------
    s = slide(prs); n += 1
    header(s, "Machine learning", "The detector: fine-tuning YOLO11")
    panel(s, 0.62, 2.0, 6.0, 3.45)
    p = s.shapes[-1]
    panel_text(p, [
        ("TRANSFER LEARNING", 13, ACC, True),
        ("", 6, INK, False),
        ("Base: YOLO11s — pre-trained on COCO (80 classes)", 15, INK, True),
        ("Re-trained on basketball footage → 3 custom classes:", 14, INK, False),
        ("ball   ·   basket   ·   person", 16, GOLD, True),
        ("", 6, INK, False),
        ("Why fine-tune? COCO has no 'basket' class and detects a small, fast ball poorly. A few epochs fixed both — the pre-trained backbone already knows generic visual features.", 13.5, MUT, False),
    ])
    panel(s, 6.9, 2.0, 5.83, 3.45, fill=RGBColor(0x12, 0x17, 0x1f))
    p = s.shapes[-1]
    panel_text(p, [
        ("TRAINING DATA", 13, ACC, True),
        ("", 4, INK, False),
        ("Roboflow Universe — “basketball-detection-dn6fg” (v4, CC BY 4.0)", 13.5, INK, True),
        ("", 4, INK, False),
        ("~7,500 labelled images", 19, INK, True),
        ("6,017 train   ·   981 val   ·   488 test", 14, MUT, False),
        ("", 6, INK, False),
        ("imgsz 512  ·  batch 8  ·  lr0 0.01  ·  optimizer auto", 13.5, INK, False),
        ("Trained on a single consumer GPU (RTX 4050)", 13.5, MUT, False, True),
    ])
    chip(s, "mAP@50   0.92", 0.62, 5.75, ACC, w=2.7)
    chip(s, "Precision   0.94", 3.55, 5.75, AWAY, w=2.7)
    chip(s, "Recall   0.86", 6.48, 5.75, GOLD, w=2.7)
    txt(s, "Validation metrics on 981 held-out images — strong accuracy after only a few epochs.",
        0.62, 6.35, 11.8, 0.4, 13, MUT, italic=True)
    notes(s, """
This is the machine-learning core. The object detector is a YOLO11-small model. We did NOT
train it from scratch — we used transfer learning: we started from weights already
pre-trained on COCO, which has eighty everyday object classes, and fine-tuned it on
basketball data with just three classes we care about — ball, basket and person. Why
fine-tune at all? Because COCO has no 'basket' or hoop class, and its generic 'sports ball'
detector is weak on a small, fast-moving basketball. The training data is a public Roboflow
Universe dataset — 'basketball-detection', about seven and a half thousand labelled images,
split into train, validation and test. We trained at 512-pixel input on a single consumer
RTX 4050 GPU. And here's the payoff of transfer learning: because the backbone already knew
generic visual features, it converged in only a few epochs and still reached a mean average
precision at fifty percent IoU of about zero-point-nine-two, with precision around
ninety-four percent and recall around eighty-six. Those are the validation numbers on
nearly a thousand held-out images.
""")
    footer(s, n, TOTAL)

    # 7 — ML: MODELS & METHODS AT A GLANCE ---------------------------------
    s = slide(prs); n += 1
    header(s, "The toolkit", "Models & methods")
    methods = [
        ("DETECTION  ·  YOLO11s (fine-tuned)", "Single-stage detector — finds the ball, basket and every player in each frame, in real time.", HOME),
        ("TRACKING  ·  ByteTrack", "Multi-object tracker — assigns and persists a track ID per player across frames (we aggregate by jersey to beat fragmentation).", AWAY),
        ("RECOGNITION  ·  PaddleOCR (GPU)", "Optical character recognition — reads jersey numbers and the on-screen scoreboard digits.", GOLD),
        ("MOTION  ·  Kalman filter", "Constant-acceleration state estimator — smooths the ball trajectory and fills gaps where detection drops.", ACC),
    ]
    y = 2.0
    for k, body, c in methods:
        p = panel(s, 0.62, y, 12.1, 1.16)
        barl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(y), Inches(0.1), Inches(1.16))
        _solid(barl, c); barl.line.fill.background(); barl.shadow.inherit = False
        panel_text(p, [(k, 14, c, True), (body, 14.5, INK, False)], pad=0.18)
        y += 1.28
    notes(s, """
Zooming out, here's the full toolkit — four learned or model-based components working
together. Detection is the YOLO model we just discussed. Tracking is ByteTrack, which gives
each player a persistent ID across frames — though camera cuts fragment those IDs, which is
why we aggregate statistics by jersey number rather than by raw track. Recognition is
PaddleOCR running on the GPU, used twice — for jersey numbers and for the scoreboard.
And motion is a Kalman filter, a classic constant-acceleration state estimator, that smooths
the ball's path and predicts through the frames where the detector loses it. So it's not a
single model — it's a pipeline of complementary techniques, each chosen for a specific
sub-problem.
""")
    footer(s, n, TOTAL)

    # 8 — KEY INSIGHT: scoreboard = ground truth ---------------------------
    s = slide(prs); n += 1
    header(s, "Key insight", "In 3×3, the scoreboard IS the ground truth", color=GOLD)
    panel(s, 0.62, 2.0, 6.0, 4.1, fill=RGBColor(0x1f, 0x1a, 0x12), line=RGBColor(0x5a, 0x3a, 0x1a))
    p = s.shapes[-1]
    panel_text(p, [
        ("Why detecting makes/misses from the ball is hard", 15, GOLD, True),
        ("Noisy ball track, occlusions, ad-logo false positives.", 14, INK, False),
        ("", 8, INK, False),
        ("So we read the scoreboard instead", 15, GOLD, True),
        ("+1  =  a make from INSIDE the arc (1 pt)", 16, INK, True),
        ("+2  =  a make from OUTSIDE the arc (2 pt)", 16, INK, True),
        ("…each with the exact game-clock time.", 14, MUT, False, True),
    ])
    panel(s, 6.9, 2.0, 5.8, 4.1)
    p = s.shapes[-1]
    panel_text(p, [
        ("RESULT", 13, ACC, True),
        ("", 6, INK, False),
        ("12 made shots recovered", 22, INK, True),
        ("9 × 1-pointer   +   3 × 2-pointer", 16, MUT, False),
        ("", 8, INK, False),
        ("Reconciles EXACTLY to the final score", 16, ACC, True),
        ("Home 11 – 10 Away  (incl. 3–3 pre-footage baseline)", 14, MUT, False),
        ("", 8, INK, False),
        ("Every shot carries its value and its clock time — far more reliable than guessing makes from the ball.", 14, INK, False),
    ])
    notes(s, """
This is the idea I'm most proud of. The obvious way to detect a made shot is to track the
ball into the hoop — but on this footage the ball track is noisy, players occlude it, and
circular ad logos cause false positives. So we flipped the problem. In 3x3 basketball the
scoreboard encodes every make for us: a plus-one is a basket from inside the arc, a plus-two
is from outside, and it changes at the exact moment of the score. So we just read the
scoreboard twice a second and reconstruct the game. The result on the right: twelve made
shots, nine ones and three twos, and crucially it reconciles exactly to the final score of
eleven-ten. When your reconstruction matches the official score perfectly, you know your
shot data is solid. The three-three is the score that was already on the board when the clip
starts, forty seconds in.
""")
    footer(s, n, TOTAL)

    # 9 — RESULT: game flow -------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Results", "Game flow, reconstructed automatically")
    pic_fit(s, R / "chart_gameflow.png", 0.62, 1.95, 8.6, 4.6)
    bullets(s, [
        ("Score progression for both teams, rebuilt purely from scoreboard OCR.", INK, False),
        ("Every dot is a made shot — large = 2-pointer, small = 1-pointer.", INK, False),
        ("Final: Home 11 – 10 Away.", INK, True),
    ], 9.45, 2.4, 3.4, 3.5, size=15, gap=12)
    notes(s, """
Now the results, all generated from that one video. This is the game flow — the score of
both teams over time, rebuilt entirely from reading the scoreboard. Each dot is a made shot,
with the big dots being two-pointers. You can read the whole story of the game off this
chart: who led, the runs, the close finish at eleven-ten. A coach or player gets this
automatically, with no tagging.
""")
    footer(s, n, TOTAL)

    # 10 — RESULT: shot distribution map ------------------------------------
    s = slide(prs); n += 1
    header(s, "Results", "Shot-distribution map")
    pic_fit(s, R / "chart_shotchart.png", 0.62, 1.95, 7.6, 4.7)
    panel(s, 8.45, 2.1, 4.3, 4.0)
    p = s.shapes[-1]
    panel_text(p, [
        ("HONEST BY DESIGN", 13, ACC, True),
        ("", 6, INK, False),
        ("Zone is real data", 16, INK, True),
        ("Inside vs outside the arc comes straight from the scoreboard (+1 / +2) — exact and per-team.", 14, INK, False),
        ("", 8, INK, False),
        ("Position within a zone is illustrative", 16, GOLD, True),
        ("The broadcast camera pans & zooms, so we don't fake pin-point xy — we show the zone the data truly supports.", 14, INK, False),
    ])
    notes(s, """
This is the flagship deliverable from the product spec: the shot-distribution map. Each
marker is a made shot on a 3x3 half-court, coloured by team and labelled by value. I want to
be precise about what's real here, because it matters. The zone — inside versus outside the
arc — is real, exact data: it comes directly from whether the scoreboard ticked up by one or
two. The exact spot within a zone is illustrative, because the broadcast camera constantly
pans and zooms, so we can't recover true court coordinates. Rather than invent fake precise
positions, we show the fidelity the footage actually supports and we say so on the slide.
That honesty is a feature, not an apology.
""")
    footer(s, n, TOTAL)

    # 11 — RESULT: team scoring / shot mix ----------------------------------
    s = slide(prs); n += 1
    header(s, "Results", "Team scoring & shot mix")
    pic_fit(s, R / "chart_shotmix.png", 0.62, 1.95, 5.4, 4.6)
    panel(s, 6.4, 2.1, 6.35, 4.2)
    p = s.shapes[-1]
    panel_text(p, [
        ("WHAT THE NUMBERS SAY", 13, ACC, True),
        ("", 6, INK, False),
        ("Home  ·  7 makes  =  6×1pt + 1×2pt  →  8 pts", 16, HOME, True),
        ("Away  ·  5 makes  =  3×1pt + 2×2pt  →  7 pts", 16, AWAY, True),
        ("", 8, INK, False),
        ("Away leaned on the two-point shot", 16, INK, True),
        ("40% of Away's makes were 2-pointers vs only 14% for Home — a real, data-backed difference in shot selection.", 14, INK, False),
        ("", 6, INK, False),
        ("(Points scored during the analyzed footage; final score 11–10 adds the 3–3 baseline.)", 12, MUT, False, True),
    ])
    notes(s, """
The same shot data, turned into team insight. Home scored on seven makes for eight points;
Away on five makes for seven points. But look at the shot mix: forty percent of Away's
makes were two-pointers, versus only fourteen percent for Home. That's a genuine,
data-backed difference in strategy — Away leaned on the outside shot. This is exactly the
kind of insight amateur teams never normally get. One note for clarity: these are the points
scored inside the footage; the eleven-ten final adds the three-three that was already on the
board when the clip began.
""")
    footer(s, n, TOTAL)

    # 12 — PLAYER CARDS -----------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Results", "Player identity cards")
    thumbs = [(5, "HIGH", ACC), (22, "HIGH", ACC), (12, "HIGH", ACC), (2, "MED", GOLD)]
    x = 0.62; w = 2.55; gap = 0.2
    for i, (j, tier, c) in enumerate(thumbs):
        lx = x + i * (w + gap)
        pic_fit(s, R / "cards" / f"thumb_{j}.png", lx, 2.0, w, 3.1)
        chip(s, f"#{j} · {tier}", lx + 0.15, 4.7, c, w=1.5)
    bullets(s, [
        ("Each player is identified by jersey-number OCR + tracking, with a real thumbnail and a confidence tag.", INK, False),
        ("Court-presence stats we can attribute reliably: time tracked, OCR reads, tracks merged.", INK, False),
        ("In the live dashboard the cards are interactive — click one to reveal its stats.", ACC, True),
    ], 0.62, 5.3, 12.1, 1.6, size=15, gap=8)
    notes(s, """
Identity. For each strongly-seen player the system builds a card: a real thumbnail cropped
from the video, the jersey number, a confidence tag, and the court-presence stats we can
attribute reliably — how long they were tracked, how many times we read their number, and so
on. In the live dashboard these cards are interactive: you click a player to reveal their
stats, rather than dumping everything at once. It feels like a modern sports app, and it's
all driven by the OCR and tracking under the hood.
""")
    footer(s, n, TOTAL)

    # 13 — HIGHLIGHTS -------------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Results", "Auto-generated highlights")
    hts = sorted((R / "highlights" / "thumbs").glob("shot_*.png"))[:4]
    x = 0.62; w = 2.95; gap = 0.18
    for i, ht in enumerate(hts):
        pic_fit(s, ht, x + i * (w + gap), 2.1, w, 1.95)
    bullets(s, [
        ("One clip per made shot, cut straight from the broadcast using the scoreboard timestamps (6 s before → 3 s after).", INK, False),
        ("Burnt-in, team-coloured captions; 12 clips plus a stitched full-game reel.", INK, False),
        ("Zero manual editing — the data decides where every clip starts and ends.", ACC, True),
    ], 0.62, 4.4, 12.1, 2.2, size=16, gap=12)
    notes(s, """
Because we know the exact clock time of every made shot, we can auto-edit. The system cuts
one highlight clip per make — six seconds before to three after — burns in a team-coloured
caption, and stitches them into a full-game reel. Twelve clips, one reel, no human editor.
The data itself decides where every clip begins and ends. This is the kind of shareable
content that gets players to actually come back and use the product.
""")
    footer(s, n, TOTAL)

    # 14 — LIMITATIONS ------------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Engineering honesty", "What this footage does NOT support", color=GOLD)
    lims = [
        ("Pin-point court coordinates", "The broadcast camera pans & zooms continuously — there is no single court-to-image mapping, so no true heat map. We report shots by zone instead."),
        ("Per-player shot attribution", "Tying each make to a jersey needs a clean ball-to-rim track AND a correct number read; validation showed neither holds reliably here."),
        ("Exact jersey numbers", "Numbers are only ~20 px tall — resolution-capped. We present the roster as 'numbers seen', weighted by confidence."),
        ("Field-goal % / eFG%", "Needs missed attempts; misses aren't detected reliably, so we report made-shot distribution only."),
    ]
    y = 2.0
    for k, body in lims:
        p = panel(s, 0.62, y, 12.1, 1.04, fill=RGBColor(0x1f, 0x1a, 0x12), line=RGBColor(0x5a, 0x3a, 0x1a))
        panel_text(p, [(k, 15, GOLD, True), (body, 13.5, INK, False)], pad=0.16)
        y += 1.14
    notes(s, """
I want to spend a moment here because this slide is what separates a demo from real
engineering. We were rigorous about what the footage does NOT support. We cannot give
pin-point court coordinates or a true heat map, because the camera never holds still. We
cannot reliably attribute each individual shot to a specific player — we tested it, and it
failed often enough that shipping it would be dishonest. The jersey numbers themselves are
only about twenty pixels tall, so they're approximate. And we can't compute shooting
percentages because we don't detect misses. Rather than fake any of this, we report what's
solid and clearly label what's illustrative. A system that knows its own limits is one you
can trust.
""")
    footer(s, n, TOTAL)

    # 15 — ROADMAP ----------------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Roadmap", "What better input unlocks")
    panel(s, 0.62, 2.0, 6.0, 4.1, fill=RGBColor(0x12, 0x1a, 0x14), line=RGBColor(0x2f, 0x6b, 0x3f))
    p = s.shapes[-1]
    panel_text(p, [
        ("TWO INDEPENDENT NEEDS", 13, ACC, True),
        ("", 6, INK, False),
        ("1 · Higher-resolution source", 16, INK, True),
        ("→ reliable jersey numbers (player identity).", 14, MUT, False),
        ("", 6, INK, False),
        ("2 · A fixed, calibrated court camera", 16, INK, True),
        ("→ true shot positions, heat maps & per-player attribution.", 14, MUT, False),
        ("", 8, INK, False),
        ("These are separate problems — resolution alone does not unlock the court-space metrics.", 14, ACC, True),
    ])
    panel(s, 6.9, 2.0, 5.8, 4.1)
    p = s.shapes[-1]
    panel_text(p, [
        ("THEN THE FULL PRD METRICS OPEN UP", 13, ACC, True),
        ("", 6, INK, False),
        ("• True shot-distribution map (real x,y)", 15, INK, False),
        ("• Effective Field-Goal %", 15, INK, False),
        ("• Rebound Control Rate (ORB / DRB)", 15, INK, False),
        ("• Catch & Shoot vs Off-the-Dribble", 15, INK, False),
        ("• Per-game growth curves (multi-game)", 15, INK, False),
        ("", 8, INK, False),
        ("The pipeline is already built to produce these — it's gated on input quality, not on the software.", 14, MUT, False, True),
    ])
    notes(s, """
So what's next, and what would unlock the rest of the product spec? Two separate things, and
it's important not to conflate them. One: a higher-resolution source fixes the jersey
numbers — that's purely an identity problem. Two: a fixed, calibrated court camera — one that
doesn't pan or zoom — unlocks true shot positions, heat maps and per-player attribution.
That's a geometry problem. Resolution alone does not solve the court-space metrics; you need
the stable camera too. With those two inputs, the full set of spec metrics opens up:
effective field-goal percentage, rebound control rate, catch-and-shoot versus off-the-dribble,
growth curves across games. And the key point — the software is already built to produce
these. We're gated on input quality, not on engineering.
""")
    footer(s, n, TOTAL)

    # 16 — CLOSING ----------------------------------------------------------
    s = slide(prs); n += 1
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    _solid(band, ACC); band.line.fill.background(); band.shadow.inherit = False
    txt(s, "SUMMARY", 0.9, 1.5, 11, 0.5, 14, ACC, bold=True)
    txt(s, "A working pipeline that turns raw 3×3 footage\ninto a shareable analytics product.",
        0.86, 1.95, 11.6, 1.5, 32, INK, bold=True)
    bullets(s, [
        ("12 made shots recovered, reconciled exactly to the final score.", INK, False),
        ("Game-flow, shot map, team analytics, player cards & 12 auto-cut highlights.", INK, False),
        ("Honest about its limits, with a clear path to the full metric set.", INK, False),
    ], 0.9, 3.9, 11.5, 2.0, size=18, gap=12)
    txt(s, "Live demo:  runs720/dashboard.html   ·   Thank you", 0.9, 6.2, 11.5, 0.5, 16, ACC, bold=True)
    notes(s, """
To wrap up. We set out to turn one raw amateur game video into the kind of analytics only
pros get — and we have a working pipeline that does it. We recover every made shot and
reconcile exactly to the final score; we generate a game-flow view, a shot map, team
analytics, interactive player cards and a dozen auto-cut highlights. And we did it honestly:
the system knows what the footage supports and what it doesn't, with a clear path to the full
metric set once the input quality improves. The best part is the live dashboard — I'd love
to show it to you now. Thank you.
""")
    footer(s, n, TOTAL)

    # 17 — CODE & RESOURCES -------------------------------------------------
    s = slide(prs); n += 1
    header(s, "Open source", "Code & resources")
    panel(s, 0.62, 2.0, 7.4, 4.1)
    p = s.shapes[-1]
    panel_text(p, [
        ("REPOSITORY", 13, ACC, True),
        (repo_url, 20, INK, True),
        ("", 8, INK, False),
        ("What's inside", 13, MUT, True),
        ("• Full pipeline source (detection → identity → analysis → output)", 14, INK, False),
        ("• Interactive HTML dashboard generator", 14, INK, False),
        ("• Highlight auto-editor + this presentation generator", 14, INK, False),
        ("• README with setup & usage", 14, INK, False),
        ("", 8, INK, False),
        ("MIT-style, reproducible — clone, install, run.", 14, ACC, True),
    ])
    panel(s, 8.25, 2.0, 4.5, 4.1, fill=RGBColor(0x12, 0x17, 0x1f))
    p = s.shapes[-1]
    panel_text(p, [
        ("RUN IT YOURSELF", 13, ACC, True),
        ("", 8, INK, False),
        ("pip install -r requirements.txt", 13, INK, True),
        ("python -m src.output.build_dashboard", 13, MUT, True),
        ("python -m src.output.build_deck", 13, MUT, True),
        ("", 8, INK, False),
        ("Built with", 13, MUT, True),
        ("Python · YOLO (Ultralytics) · PaddleOCR · OpenCV · pandas · matplotlib · ffmpeg",
         14, INK, False),
    ])
    txt(s, "Thank you — questions welcome.", 0.62, 6.45, 11.5, 0.5, 16, ACC, bold=True)
    notes(s, """
Finally, this is open. The complete codebase is on GitHub at the address shown — the full
pipeline, the dashboard generator, the highlight auto-editor, even the generator for these
very slides, plus a README with setup and usage. It's reproducible: clone it, install the
requirements, point it at a video and run. Everything you've seen today came out of that
repository, from one raw video file. Thank you very much — I'm happy to take any questions.
""")
    footer(s, n, TOTAL)

    prs.save(str(out))
    print(f"Wrote {out}  ({out.stat().st_size/1024:.0f} KB, {len(prs.slides._sldIdLst)} slides)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--runs", type=Path, default=Path("runs720"))
    ap.add_argument("--out", type=Path, default=Path("Basketball_AI_Presentation.pptx"))
    ap.add_argument("--repo-url", default="github.com/marcoveron/basketball-ai")
    args = ap.parse_args()
    build(args.runs, args.out, args.repo_url)


if __name__ == "__main__":
    main()
